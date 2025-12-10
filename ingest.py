# ingest.py (UPDATED with PDF Image OCR + DOCX Image OCR + PPTX Image OCR)

import os
import shutil
import io
from docx import Document as DocxDocument
from PyPDF2 import PdfReader
import pdfplumber
from pptx import Presentation
from PIL import Image
import pytesseract

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# ------------------------------------------------------------
# TEXT SPLITTER
# ------------------------------------------------------------
def text_splitting_recursive(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        length_function=len,
    )
    return splitter.split_text(text)

# ------------------------------------------------------------
# PDF EXTRACTION (TEXT + IMAGE OCR)
# ------------------------------------------------------------
def extract_text_from_pdf(file_path):
    text = ""

    # 1️⃣ Try extracting actual PDF text
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t
    except Exception as e:
        print(f"⚠️ PyPDF2 failed on {file_path}: {e}")

    # 2️⃣ If real text empty → try pdfplumber
    if not text.strip():
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text() or ""
                    text += extracted
        except Exception as e:
            print(f"⚠️ pdfplumber failed on {file_path}: {e}")

    # 3️⃣ OCR on images inside PDF pages
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                images = page.images
                for img_info in images:
                    try:
                        x = img_info.get("x0")
                        y = img_info.get("top")
                        w = img_info.get("width")
                        h = img_info.get("height")

                        cropped = page.crop((x, y, x + w, y + h)).to_image(resolution=300)
                        pil_img = cropped.original
                        ocr_text = pytesseract.image_to_string(pil_img)
                        text += "" + ocr_text

                    except Exception:
                        pass
    except Exception:
        pass

    return text

# ------------------------------------------------------------
# DOCX EXTRACTION (TEXT + IMAGE OCR)
# ------------------------------------------------------------
def extract_text_from_docx(file_path):
    doc = DocxDocument(file_path)
    all_text = []

    # 1️⃣ Extract paragraphs
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            all_text.append(paragraph.text.strip())

    # 2️⃣ Extract tables
    for table in doc.tables:
        for row in table.rows:
            cell_texts = []
            for cell in row.cells:
                if cell.text.strip():
                    cell_texts.append(cell.text.strip())
            if cell_texts:
                all_text.append(" | ".join(cell_texts))

    # 3️⃣ OCR on embedded images inside DOCX (images stored separately)
    doc_folder = os.path.join(os.path.dirname(file_path), "_docx_temp_extract")
    if os.path.exists(doc_folder):
        shutil.rmtree(doc_folder)
    os.makedirs(doc_folder, exist_ok=True)

    for rel in doc.part._rels:
        try:
            target = doc.part._rels[rel].target_ref
            if "image" in target:
                image_part = doc.part._rels[rel].target_part
                img_bytes = image_part.blob

                img_path = os.path.join(doc_folder, "img.png")
                with open(img_path, "wb") as f:
                    f.write(img_bytes)

                img = Image.open(img_path)
                ocr_text = pytesseract.image_to_string(img)
                if ocr_text.strip():
                    all_text.append(ocr_text.strip())
        except Exception:
            pass

    return "".join(all_text)

# ------------------------------------------------------------
# PPTX EXTRACTION (TEXT + IMAGE OCR)
# ------------------------------------------------------------
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    all_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())

    return "".join(all_text)

def extract_images_from_pptx(file_path):
    prs = Presentation(file_path)
    images = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                img = shape.image
                try:
                    pil_img = Image.open(io.BytesIO(img.blob)).convert("RGB")
                    images.append(pil_img)
                except:
                    pass

    return images

# ------------------------------------------------------------
# STANDALONE IMAGE OCR
# ------------------------------------------------------------
IMAGE_EXTENSIONS = [".png", ".jpg", ".jpeg", ".webp"]

def extract_text_from_image(file_path):
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        print(f"⚠️ OCR failed for {file_path}: {e}")
        return ""

# ------------------------------------------------------------
# INDEXING MAIN FUNCTION
# ------------------------------------------------------------
def index_files(file_paths, faiss_dir, progress_callback=None):

    # CLEAN OLD INDEX
    if os.path.exists(faiss_dir):
        try:
            shutil.rmtree(faiss_dir)
        except Exception:
            pass
    os.makedirs(faiss_dir, exist_ok=True)

    documents = []
    total_files = len(file_paths)

    # folder for extracted images
    extracted_folder = os.path.join(faiss_dir, "extracted_images")
    os.makedirs(extracted_folder, exist_ok=True)

    for i, file in enumerate(file_paths, start=1):

        ext = os.path.splitext(file)[1].lower()

        if progress_callback:
            progress_callback(
                i / (total_files * 2),
                f"📄 Processing file {i}/{total_files}: {os.path.basename(file)}"
            )

        text = ""
        image_paths = []   # for linking images to metadata

        # -----------------------------
        # PDF
        # -----------------------------
        if ext == ".pdf":
            text = extract_text_from_pdf(file)

        # -----------------------------
        # DOCX
        # -----------------------------
        elif ext == ".docx":
            text = extract_text_from_docx(file)

        # -----------------------------
        # PPTX
        # -----------------------------
        elif ext == ".pptx":
            text = extract_text_from_pptx(file)

            # OCR IMAGES in PPTX
            pptx_images = extract_images_from_pptx(file)
            for idx, img in enumerate(pptx_images):
                img_path = os.path.join(extracted_folder, f"{os.path.basename(file)}_img{idx}.jpg")
                img.save(img_path)
                image_paths.append(img_path)

                # OCR text
                try:
                    text += "\n" + pytesseract.image_to_string(img)
                except:
                    pass

        # -----------------------------
        # IMAGE FILES (PNG/JPG)
        # -----------------------------
        elif ext in IMAGE_EXTENSIONS:
            text = extract_text_from_image(file)

            # save image for retrieval
            img_output_path = os.path.join(extracted_folder, os.path.basename(file))
            shutil.copy(file, img_output_path)
            image_paths.append(img_output_path)

        # -----------------------------
        else:
            print(f"⚠️ Unsupported file skipped: {file}")
            continue

        if not text.strip():
            print(f"⚠️ No text extracted from {os.path.basename(file)}")
            continue

        # Split into chunks
        chunks = text_splitting_recursive(text)

        for chunk in chunks:
            doc_metadata = {"source": file}

            # attach image metadata if exists
            if image_paths:
                doc_metadata["image_paths"] = image_paths

            documents.append(
                Document(page_content=chunk, metadata=doc_metadata)
            )

    # -----------------------------
    # No documents → stop
    # -----------------------------
    if not documents:
        if progress_callback:
            progress_callback(1.0, "⚠️ No valid files to index.")
        return {"path": None, "total_chunks": 0, "files_indexed": 0}

    # -----------------------------
    # Embed & Save FAISS
    # -----------------------------
    if progress_callback:
        progress_callback(0.9, "🧠 Generating embeddings...")

    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    db = FAISS.from_documents(documents, embeddings)

    save_path = os.path.abspath(faiss_dir)
    db.save_local(save_path)

    if progress_callback:
        progress_callback(1.0, "✅ Indexing complete!")

    print(f"✅ Indexed {len(documents)} chunks from {len(file_paths)} files.")

    return {
        "path": save_path,
        "total_chunks": len(documents),
        "files_indexed": total_files,
    }
